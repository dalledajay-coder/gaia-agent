"""File handling tools for processing GAIA task attachments."""

import os
import base64
import subprocess
from claude_agent_sdk import tool
from typing import Any


@tool(
    "read_file",
    "Read the contents of a file. Supports text, CSV, JSON, Excel, PDF, DOCX, and provides info for binary files like images/audio. For images, describes them. For audio, attempts transcription.",
    {"file_path": str},
)
async def read_file(args: dict[str, Any]) -> dict[str, Any]:
    file_path = args["file_path"]

    if not os.path.exists(file_path):
        return {"content": [{"type": "text", "text": f"File not found: {file_path}"}]}

    ext = os.path.splitext(file_path)[1].lower()

    try:
        # Text-based files
        text_exts = {
            ".txt", ".csv", ".json", ".md", ".py", ".js", ".html", ".xml",
            ".yaml", ".yml", ".tsv", ".log", ".cfg", ".ini", ".toml",
            ".tex", ".rst", ".r", ".sql", ".sh", ".bat", ".css",
            ".jsonl", ".ndjson", ".geojson",
        }
        if ext in text_exts:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            if len(content) > 20000:
                content = content[:20000] + "\n... [truncated]"
            return {"content": [{"type": "text", "text": f"File: {file_path}\n\n{content}"}]}

        elif ext in (".xlsx", ".xls"):
            return await _read_excel(file_path)

        elif ext == ".pdf":
            return await _read_pdf(file_path)

        elif ext in (".docx", ".doc"):
            return await _read_docx(file_path)

        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"):
            return await _describe_image(file_path)

        elif ext in (".mp3", ".wav", ".ogg", ".flac", ".m4a", ".mp4", ".webm"):
            return await _process_audio(file_path)

        elif ext == ".pptx":
            return await _read_pptx(file_path)

        else:
            # Try reading as text first
            try:
                with open(file_path, "r", encoding="utf-8", errors="strict") as f:
                    content = f.read(20000)
                return {"content": [{"type": "text", "text": f"File: {file_path}\n\n{content}"}]}
            except (UnicodeDecodeError, ValueError):
                size = os.path.getsize(file_path)
                return {"content": [{"type": "text", "text": f"Binary file: {file_path} ({size} bytes, type: {ext}). Use execute_python to process this file with appropriate libraries."}]}

    except Exception as e:
        return {"content": [{"type": "text", "text": f"Error reading {file_path}: {str(e)}"}]}


async def _read_excel(file_path: str) -> dict[str, Any]:
    """Read Excel files using openpyxl."""
    try:
        result = subprocess.run(
            ["python3", "-c", f"""
import sys
try:
    import openpyxl
    wb = openpyxl.load_workbook('{file_path}', data_only=True)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        print(f"=== Sheet: {{sheet_name}} ===")
        for row in ws.iter_rows(values_only=True):
            print('\\t'.join(str(c) if c is not None else '' for c in row))
except ImportError:
    try:
        import pandas as pd
        for sheet_name in pd.ExcelFile('{file_path}').sheet_names:
            df = pd.read_excel('{file_path}', sheet_name=sheet_name)
            print(f"=== Sheet: {{sheet_name}} ===")
            print(df.to_string())
    except ImportError:
        print("Neither openpyxl nor pandas available")
"""],
            capture_output=True, text=True, timeout=30
        )
        if result.stdout:
            text = result.stdout[:20000]
            return {"content": [{"type": "text", "text": f"Excel file: {file_path}\n\n{text}"}]}
        return {"content": [{"type": "text", "text": f"Could not read Excel file: {result.stderr}"}]}
    except Exception as e:
        return {"content": [{"type": "text", "text": f"Excel read error: {str(e)}"}]}


async def _read_pdf(file_path: str) -> dict[str, Any]:
    """Read PDF files using multiple methods."""
    try:
        result = subprocess.run(
            ["python3", "-c", f"""
import sys
text_parts = []

# Method 1: pdftotext
try:
    import subprocess as sp
    r = sp.run(['pdftotext', '{file_path}', '-'], capture_output=True, text=True, timeout=30)
    if r.stdout and len(r.stdout.strip()) > 50:
        text_parts.append(r.stdout)
except Exception:
    pass

# Method 2: PyPDF2
if not text_parts:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader('{file_path}')
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text_parts.append(t)
    except Exception:
        pass

if text_parts:
    print('\\n'.join(text_parts))
else:
    # Return page count and metadata
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader('{file_path}')
        print(f"PDF with {{len(reader.pages)}} pages. Text extraction failed - may contain images only.")
    except Exception:
        print("Could not read PDF")
"""],
            capture_output=True, text=True, timeout=60
        )
        text = result.stdout if result.stdout else f"Could not extract text: {result.stderr}"
        if len(text) > 20000:
            text = text[:20000] + "\n... [truncated]"
        return {"content": [{"type": "text", "text": f"PDF file: {file_path}\n\n{text}"}]}
    except Exception as e:
        return {"content": [{"type": "text", "text": f"PDF read error: {str(e)}"}]}


async def _read_docx(file_path: str) -> dict[str, Any]:
    """Read DOCX files."""
    try:
        result = subprocess.run(
            ["python3", "-c", f"""
try:
    from docx import Document
    doc = Document('{file_path}')
    for para in doc.paragraphs:
        print(para.text)
    for table in doc.tables:
        for row in table.rows:
            print('\\t'.join(cell.text for cell in row.cells))
except ImportError:
    # Fallback: use zipfile to extract text
    import zipfile
    import xml.etree.ElementTree as ET
    with zipfile.ZipFile('{file_path}') as z:
        with z.open('word/document.xml') as f:
            tree = ET.parse(f)
            for elem in tree.iter():
                if elem.text:
                    print(elem.text, end=' ')
"""],
            capture_output=True, text=True, timeout=30
        )
        text = result.stdout if result.stdout else "Could not read DOCX file"
        if len(text) > 20000:
            text = text[:20000] + "\n... [truncated]"
        return {"content": [{"type": "text", "text": f"DOCX file: {file_path}\n\n{text}"}]}
    except Exception as e:
        return {"content": [{"type": "text", "text": f"DOCX read error: {str(e)}"}]}


async def _describe_image(file_path: str) -> dict[str, Any]:
    """Provide information about an image file."""
    size = os.path.getsize(file_path)
    ext = os.path.splitext(file_path)[1].lower()

    info = f"Image file: {file_path} ({size} bytes, format: {ext})\n\n"
    info += "To analyze this image, use execute_python with PIL/Pillow to inspect properties, "
    info += "or use the Bash tool to run image analysis commands.\n"

    # Try to get image dimensions
    try:
        result = subprocess.run(
            ["python3", "-c", f"""
from PIL import Image
img = Image.open('{file_path}')
print(f"Dimensions: {{img.size[0]}}x{{img.size[1]}}")
print(f"Mode: {{img.mode}}")
print(f"Format: {{img.format}}")

# If image has text, try OCR
try:
    import pytesseract
    text = pytesseract.image_to_string(img)
    if text.strip():
        print(f"\\nOCR Text:\\n{{text}}")
except ImportError:
    print("\\nOCR not available (install pytesseract)")
"""],
            capture_output=True, text=True, timeout=30
        )
        if result.stdout:
            info += "\n" + result.stdout
    except Exception:
        pass

    return {"content": [{"type": "text", "text": info}]}


async def _process_audio(file_path: str) -> dict[str, Any]:
    """Process audio files - provide info and attempt transcription."""
    size = os.path.getsize(file_path)
    ext = os.path.splitext(file_path)[1].lower()

    info = f"Audio/video file: {file_path} ({size} bytes, format: {ext})\n\n"
    info += "To process this audio, use execute_python with appropriate libraries.\n"
    info += "You can use the Bash tool to run ffmpeg commands for conversion.\n"

    # Try to get duration
    try:
        result = subprocess.run(
            ["ffprobe", "-i", file_path, "-show_entries", "format=duration",
             "-v", "quiet", "-of", "csv=p=0"],
            capture_output=True, text=True, timeout=10
        )
        if result.stdout.strip():
            info += f"\nDuration: {float(result.stdout.strip()):.1f} seconds"
    except Exception:
        pass

    return {"content": [{"type": "text", "text": info}]}


async def _read_pptx(file_path: str) -> dict[str, Any]:
    """Read PowerPoint files."""
    try:
        result = subprocess.run(
            ["python3", "-c", f"""
try:
    from pptx import Presentation
    prs = Presentation('{file_path}')
    for i, slide in enumerate(prs.slides, 1):
        print(f"=== Slide {{i}} ===")
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                print(shape.text)
except ImportError:
    import zipfile
    import xml.etree.ElementTree as ET
    with zipfile.ZipFile('{file_path}') as z:
        for name in sorted(z.namelist()):
            if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
                with z.open(name) as f:
                    tree = ET.parse(f)
                    for elem in tree.iter():
                        if elem.text and elem.text.strip():
                            print(elem.text.strip())
"""],
            capture_output=True, text=True, timeout=30
        )
        text = result.stdout if result.stdout else "Could not read PPTX file"
        return {"content": [{"type": "text", "text": f"PowerPoint file: {file_path}\n\n{text[:20000]}"}]}
    except Exception as e:
        return {"content": [{"type": "text", "text": f"PPTX read error: {str(e)}"}]}
